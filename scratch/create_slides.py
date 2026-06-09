import os
import re
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from svglib.svglib import svg2rlg
from reportlab.graphics import renderPM
from PIL import Image

# ---------------------------------------------------------------------------
# SVG to PNG Conversion with CSS Variable Resolution & Background Adjustment
# ---------------------------------------------------------------------------
def convert_svg_to_png(svg_path, png_path, target_bg=(30, 41, 59)):
    print(f"Converting {svg_path} to {png_path}...")
    # Read SVG content
    with open(svg_path, 'r', encoding='utf-8') as f:
        svg_text = f.read()
    
    # Resolve CSS variables
    replacements = {
        'var(--_node-fill)': '#181818',
        'var(--_node-stroke)': '#CCCCCC',
        'var(--_text)': '#CCCCCC',
        'var(--_line)': '#CCCCCC',
        'var(--_arrow)': '#0078D4',
        'var(--bg)': '#1F1F1F'
    }
    for var, hex_val in replacements.items():
        svg_text = svg_text.replace(var, hex_val)
    
    # Write to a temp SVG
    temp_svg = svg_path + '.temp.svg'
    with open(temp_svg, 'w', encoding='utf-8') as f:
        f.write(svg_text)
        
    try:
        # Convert SVG to PNG
        drawing = svg2rlg(temp_svg)
        renderPM.drawToFile(drawing, png_path, fmt='PNG')
        
        # Post-process PNG: replace white background with PowerPoint card color
        img = Image.open(png_path).convert("RGBA")
        data = img.getdata()
        new_data = []
        for item in data:
            # Replace white background (and near-white anti-aliased pixels)
            if item[0] > 230 and item[1] > 230 and item[2] > 230:
                new_data.append(target_bg + (255,))
            else:
                new_data.append(item)
        img.putdata(new_data)
        img.save(png_path)
    finally:
        # Clean up temp file
        if os.path.exists(temp_svg):
            os.remove(temp_svg)
    print(f"Successfully converted {svg_path}")

# Run conversion for the 3 SVGs
convert_svg_to_png('present/3-layer grayscale.svg', 'present/3-layer grayscale.png')
convert_svg_to_png('present/3-layer rgb.svg', 'present/3-layer rgb.png')
convert_svg_to_png('present/4-layer rgb.svg', 'present/4-layer rgb.png')

# ---------------------------------------------------------------------------
# Styling and Palette (Slate Dark Dashboard Theme)
# ---------------------------------------------------------------------------
BG_COLOR = RGBColor(15, 23, 42)       # Slate 900 (Deep dark background)
CARD_BG_COLOR = RGBColor(30, 41, 59)   # Slate 800 (Card backgrounds)
CARD_BORDER_COLOR = RGBColor(71, 85, 105) # Slate 600 (Card borders)
TEXT_WHITE = RGBColor(243, 244, 246)  # Gray 100 (Title and primary text)
TEXT_MUTED = RGBColor(156, 163, 175)  # Gray 400 (Secondary text)
ACCENT_TEAL = RGBColor(45, 212, 191)  # Teal 400 (Highlight text & accents)
ACCENT_ORANGE = RGBColor(251, 146, 60) # Orange 400 (Highlight warnings & limits)

FONT_HEADING = "Arial"
FONT_BODY = "Arial"

def set_slide_background(slide):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = BG_COLOR

def add_header(slide, title_text, category_text="TINYML SMART PANTRY"):
    # Category / Tag
    tag_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.4), Inches(11.7), Inches(0.4))
    tf_tag = tag_box.text_frame
    tf_tag.word_wrap = True
    p_tag = tf_tag.paragraphs[0]
    p_tag.text = category_text.upper()
    p_tag.font.name = FONT_HEADING
    p_tag.font.size = Pt(10)
    p_tag.font.bold = True
    p_tag.font.color.rgb = ACCENT_TEAL
    
    # Main Title
    title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(11.7), Inches(0.8))
    tf_title = title_box.text_frame
    tf_title.word_wrap = True
    tf_title.margin_top = 0
    tf_title.margin_bottom = 0
    p_title = tf_title.paragraphs[0]
    p_title.text = title_text
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(28)
    p_title.font.bold = True
    p_title.font.color.rgb = TEXT_WHITE

def create_card(slide, left, top, width, height, bg_color=CARD_BG_COLOR, border_color=CARD_BORDER_COLOR):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg_color
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = Pt(1.5)
    else:
        shape.line.fill.background()
    return shape

# ---------------------------------------------------------------------------
# Presentation Initialization
# ---------------------------------------------------------------------------
prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
blank_layout = prs.slide_layouts[6]

# ===========================================================================
# SLIDE 1: Title Slide (Intro)
# ===========================================================================
slide1 = prs.slides.add_slide(blank_layout)
set_slide_background(slide1)

create_card(slide1, Inches(1.0), Inches(1.5), Inches(11.333), Inches(4.5))

main_title_box = slide1.shapes.add_textbox(Inches(1.5), Inches(2.2), Inches(10.333), Inches(3.0))
tf1 = main_title_box.text_frame
tf1.word_wrap = True

p1 = tf1.paragraphs[0]
p1.text = "TinyML Smart Pantry & Waste Reducer"
p1.font.name = FONT_HEADING
p1.font.size = Pt(40)
p1.font.bold = True
p1.font.color.rgb = TEXT_WHITE
p1.space_after = Pt(20)

p2 = tf1.add_paragraph()
p2.text = "An Engineering Journey: Physical Constraints, Optimization, and Lessons Learned"
p2.font.name = FONT_BODY
p2.font.size = Pt(18)
p2.font.color.rgb = ACCENT_TEAL
p2.space_after = Pt(40)

p3 = tf1.add_paragraph()
p3.text = "Hardware Target: Arduino Nano 33 BLE Sense + OV7675 Camera"
p3.font.name = FONT_BODY
p3.font.size = Pt(14)
p3.font.color.rgb = TEXT_MUTED

# ===========================================================================
# SLIDE 2: Core Vision & Concept
# ===========================================================================
slide2 = prs.slides.add_slide(blank_layout)
set_slide_background(slide2)
add_header(slide2, "Core Vision: Privacy-Preserving Food Tracking")

create_card(slide2, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
left_box = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
tf_left = left_box.text_frame
tf_left.word_wrap = True

p_left_title = tf_left.paragraphs[0]
p_left_title.text = "The Pantry Problem"
p_left_title.font.name = FONT_HEADING
p_left_title.font.size = Pt(20)
p_left_title.font.bold = True
p_left_title.font.color.rgb = ACCENT_TEAL
p_left_title.space_after = Pt(15)

bullets_left = [
    ("Goal", "Automatically log groceries at the pantry entrypoint and estimate expiration using USDA FoodKeeper to reduce waste."),
    ("Edge Privacy", "Perform all computer vision inference locally on-device. Images never leave the microcontroller."),
    ("Trigger Flow", "Pantry door sensor activates the camera, captures a short image, and transmits a simple event label + timestamp over BLE/Serial.")
]
for title, text in bullets_left:
    p = tf_left.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

create_card(slide2, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
right_box = slide2.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.2))
tf_right = right_box.text_frame
tf_right.word_wrap = True

p_right_title = tf_right.paragraphs[0]
p_right_title.text = "Target Hardware Constraints"
p_right_title.font.name = FONT_HEADING
p_right_title.font.size = Pt(20)
p_right_title.font.bold = True
p_right_title.font.color.rgb = ACCENT_ORANGE
p_right_title.space_after = Pt(15)

bullets_right = [
    ("Processor", "ARM Cortex-M4F microcontroller running at 64 MHz."),
    ("Flash Memory", "1 MB - plenty of space for compiling optimized neural net binaries."),
    ("SRAM Memory", "256 KB total. Crucially, Mbed OS and BLE firmware consume ~120-160 KB, leaving only 80-100 KB for model activations!"),
    ("Sensor", "OV7675 image sensor capturing low-resolution QQVGA frames.")
]
for title, text in bullets_right:
    p = tf_right.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

# ===========================================================================
# SLIDE 3: Model Iterations & Design Paths
# ===========================================================================
slide3 = prs.slides.add_slide(blank_layout)
set_slide_background(slide3)
add_header(slide3, "Model Design Iterations & Optimization")

col_width = Inches(3.6)
col_gap = Inches(0.4)
start_x = Inches(0.8)
y_pos = Inches(1.8)
height = Inches(4.8)

architectures = [
    {
        "name": "3-Layer Grayscale",
        "color": TEXT_WHITE,
        "points": [
            "Input: 96x96x1 (Grayscale)",
            "Architecture: 3 Conv2D layers with Batch Normalization & ReLU",
            "Knowledge Distillation: Student model trained to mimic a larger MobileNetV2 RGB teacher",
            "Quantization: Post-training quantized to 8-bit integer weights/activations"
        ]
    },
    {
        "name": "3-Layer RGB",
        "color": ACCENT_TEAL,
        "points": [
            "Input: 96x96x3 (RGB Color)",
            "Architecture: 3 Conv2D layers with Batch Normalization & ReLU",
            "Why Color: Added RGB channels to provide color boundaries (essential for telling foods apart)",
            "Quantization: Quantization-aware training (QAT) to maintain accuracy"
        ]
    },
    {
        "name": "4-Layer RGB",
        "color": ACCENT_ORANGE,
        "points": [
            "Input: 96x96x3 (RGB Color)",
            "Architecture: Expanded to 4 Conv2D layers for higher capacity",
            "Goal: Squeeze maximum capacity to correctly classify complex dishes",
            "Quantization: QAT to INT8 for local deployment"
        ]
    }
]

for i, arch in enumerate(architectures):
    x_pos = start_x + i * (col_width + col_gap)
    create_card(slide3, x_pos, y_pos, col_width, height)
    
    box = slide3.shapes.add_textbox(x_pos + Inches(0.25), y_pos + Inches(0.3), col_width - Inches(0.5), height - Inches(0.6))
    tf = box.text_frame
    tf.word_wrap = True
    
    p_title = tf.paragraphs[0]
    p_title.text = arch["name"]
    p_title.font.name = FONT_HEADING
    p_title.font.size = Pt(18)
    p_title.font.bold = True
    p_title.font.color.rgb = arch["color"]
    p_title.space_after = Pt(20)
    
    for pt in arch["points"]:
        p = tf.add_paragraph()
        p.text = "\u2022  " + pt
        p.font.name = FONT_BODY
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT_MUTED
        p.space_after = Pt(10)

# ===========================================================================
# SLIDE 4 (NEW): Model Architecture Comparison (SVG Pics)
# ===========================================================================
slide4_new = prs.slides.add_slide(blank_layout)
set_slide_background(slide4_new)
add_header(slide4_new, "Model Architecture Comparison")

# Sequence: 3-layer grayscale -> 3-layer rgb -> 4-layer rgb
# Position: Side-by-side corresponding to the three cards on the previous slide
# Centers: 2.6", 6.6", 10.6"
# Heights: 4.8"

# 1. 3-Layer Grayscale
create_card(slide4_new, Inches(0.8), Inches(1.8), Inches(3.6), Inches(4.8))
slide4_new.shapes.add_picture(
    'present/3-layer grayscale.png',
    left=Inches(1.51),
    top=Inches(1.8),
    width=Inches(2.18),
    height=Inches(4.8)
)

# 2. 3-Layer RGB
create_card(slide4_new, Inches(4.8), Inches(1.8), Inches(3.6), Inches(4.8))
slide4_new.shapes.add_picture(
    'present/3-layer rgb.png',
    left=Inches(5.96),
    top=Inches(1.8),
    width=Inches(1.28),
    height=Inches(4.8)
)

# 3. 4-Layer RGB
create_card(slide4_new, Inches(8.8), Inches(1.8), Inches(3.6), Inches(4.8))
slide4_new.shapes.add_picture(
    'present/4-layer rgb.png',
    left=Inches(9.96),
    top=Inches(1.8),
    width=Inches(1.28),
    height=Inches(4.8)
)

# ===========================================================================
# SLIDE 5: Results & The Memory Wall
# ===========================================================================
slide5 = prs.slides.add_slide(blank_layout)
set_slide_background(slide5)
add_header(slide5, "The Results & The Physical SRAM Wall")

left_t = Inches(0.8)
top_t = Inches(1.8)
width_t = Inches(11.7)
height_t = Inches(3.0)

table_shape = slide5.shapes.add_table(5, 5, left_t, top_t, width_t, height_t)
table = table_shape.table

table.columns[0].width = Inches(3.6)
table.columns[1].width = Inches(1.6)
table.columns[2].width = Inches(2.1)
table.columns[3].width = Inches(2.2)
table.columns[4].width = Inches(2.2)

headers = ["Model Configuration", "Input Space", "Val Accuracy (QAT)", "Binary Size (Flash)", "Peak Activations (RAM)"]
for col_idx, header in enumerate(headers):
    cell = table.cell(0, col_idx)
    cell.text = header
    cell.fill.solid()
    cell.fill.fore_color.rgb = CARD_BG_COLOR
    for p in cell.text_frame.paragraphs:
        p.font.name = FONT_HEADING
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = ACCENT_TEAL
        p.alignment = PP_ALIGN.LEFT

rows_data = [
    ["3-Layer Grayscale (Distilled)", "96x96x1", "45.6%", "~149 KB", "~305 KB"],
    ["3-Layer RGB (Float)", "96x96x3", "58.0%", "~150 KB", "Exceeds SRAM"],
    ["4-Layer RGB (Best Capacity)", "96x96x3", "72.2%", "~150 KB", "~323 KB (Exceeds SRAM)"],
    ["Nano 33 BLE Hardware Limit", "—", "—", "1 MB (Flash limit)", "80-100 KB (SRAM limit)"]
]

for row_idx, row in enumerate(rows_data):
    for col_idx, val in enumerate(row):
        cell = table.cell(row_idx + 1, col_idx)
        cell.text = val
        cell.fill.solid()
        if row_idx == 3:
            cell.fill.fore_color.rgb = RGBColor(69, 26, 26)
        else:
            cell.fill.fore_color.rgb = CARD_BG_COLOR
            
        for p in cell.text_frame.paragraphs:
            p.font.name = FONT_BODY
            p.font.size = Pt(13)
            p.font.color.rgb = TEXT_WHITE if row_idx != 3 else ACCENT_ORANGE
            if col_idx == 0 and row_idx == 3:
                p.font.bold = True
            p.alignment = PP_ALIGN.LEFT

create_card(slide5, Inches(0.8), Inches(5.1), Inches(11.7), Inches(1.5), bg_color=CARD_BG_COLOR)
warning_box = slide5.shapes.add_textbox(Inches(1.1), Inches(5.2), Inches(11.1), Inches(1.3))
tf_warn = warning_box.text_frame
tf_warn.word_wrap = True

p_warn = tf_warn.paragraphs[0]
p_warn.text = "Critical Bottleneck Analysis: RAM vs. Flash"
p_warn.font.name = FONT_HEADING
p_warn.font.size = Pt(16)
p_warn.font.bold = True
p_warn.font.color.rgb = ACCENT_ORANGE
p_warn.space_after = Pt(8)

p_warn_desc = tf_warn.add_paragraph()
p_warn_desc.text = "In TinyML, model binary size (stored in Flash) is rarely the primary constraint; activation memory (RAM) is the hard boundary. While our 4-layer RGB model compressed to 150 KB and easily fit within the 1 MB Flash, its peak activation tensors required ~323 KB. This far exceeds the ~80-100 KB available SRAM on the Arduino, causing compilation and runtime allocation failures."
p_warn_desc.font.name = FONT_BODY
p_warn_desc.font.size = Pt(13)
p_warn_desc.font.color.rgb = TEXT_MUTED

# ===========================================================================
# SLIDE 6: The Pivot & Sensor Limitations
# ===========================================================================
slide6 = prs.slides.add_slide(blank_layout)
set_slide_background(slide6)
add_header(slide6, "The Pivot & Dataset Domain Gaps")

create_card(slide6, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
pivot_box = slide6.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
tf_pivot = pivot_box.text_frame
tf_pivot.word_wrap = True

p_p_title = tf_pivot.paragraphs[0]
p_p_title.text = "The Pivot: Pizza vs. Apple Pie"
p_p_title.font.name = FONT_HEADING
p_p_title.font.size = Pt(20)
p_p_title.font.bold = True
p_p_title.font.color.rgb = ACCENT_TEAL
p_p_title.space_after = Pt(15)

bullets_pivot = [
    ("Binary Scope", "To bypass memory constraints, we reduced the problem to a simple binary task (Pizza vs. Apple Pie)."),
    ("Feature Overlap", "Both foods are round, crust-based, and baked. Distinguishing them without color (grayscale) and at 96x96 resolution remained highly unreliable."),
    ("Resolution Limits", "96x96 resolution completely washes away the fine texture detail required to identify structural features of dishes.")
]
for title, text in bullets_pivot:
    p = tf_pivot.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

create_card(slide6, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
gap_box = slide6.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.2))
tf_gap = gap_box.text_frame
tf_gap.word_wrap = True

p_g_title = tf_gap.paragraphs[0]
p_g_title.text = "Dataset Domain Mismatch"
p_g_title.font.name = FONT_HEADING
p_g_title.font.size = Pt(20)
p_g_title.font.bold = True
p_g_title.font.color.rgb = ACCENT_ORANGE
p_g_title.space_after = Pt(15)

bullets_gap = [
    ("Generalist Data", "Food-101 is web-scraped for general computer vision tasks, not industrial edge sensors."),
    ("Label Noise", "Many images contain busy backgrounds, macro shots, or people (e.g. a person taking a selfie with a cake is labeled 'cheesecake')."),
    ("Pantry Gap", "A network trained on diverse web photos cannot bridge the domain gap to top-down, fixed-lighting pantry camera captures.")
]
for title, text in bullets_gap:
    p = tf_gap.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

# ===========================================================================
# SLIDE 7: Demo — The Hardware's Perspective
# ===========================================================================
slide7 = prs.slides.add_slide(blank_layout)
set_slide_background(slide7)
add_header(slide7, "Demo: What the Hardware Actually Sees")

create_card(slide7, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8), bg_color=RGBColor(5, 5, 10))
ascii_box = slide7.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(5.2), Inches(4.4))
tf_ascii = ascii_box.text_frame
tf_ascii.word_wrap = False

p_a_title = tf_ascii.paragraphs[0]
p_a_title.text = "Arduino ASCII Serial Dump"
p_a_title.font.name = "Consolas"
p_a_title.font.size = Pt(12)
p_a_title.font.bold = True
p_a_title.font.color.rgb = ACCENT_TEAL
p_a_title.space_after = Pt(10)

ascii_art = [
    "  . . . . x x # # # # # x x . . . .  ",
    "  . . x x # # # # # # # # # x . . .  ",
    "  . x # # # # # # # # # # # # x . .  ",
    "  . x # # # # # # # # # # # # x . .  ",
    "  x # # # # # # # # # # # # # # x .  ",
    "  x # # # # # . . . . # # # # # x .  ",
    "  # # # # x . . . . . . x # # # # .  ",
    "  # # # x . . . . . . . . x # # # .  ",
    "  # # # x . . . . . . . . x # # # .  ",
    "  # # # # x . . . . . . x # # # # .  ",
    "  x # # # # # . . . . # # # # # x .  ",
    "  x # # # # # # # # # # # # # # x .  ",
    "  . x # # # # # # # # # # # # x . .  ",
    "  . x # # # # # # # # # # # # x . .  ",
    "  . . x x # # # # # # # # # x . . .  ",
    "  . . . . x x # # # # # x x . . . .  ",
]

for line in ascii_art:
    p = tf_ascii.add_paragraph()
    p.text = line
    p.font.name = "Consolas"
    p.font.size = Pt(11)
    p.font.color.rgb = RGBColor(34, 197, 94)
    p.space_after = Pt(1)

create_card(slide7, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
explain_box = slide7.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.2))
tf_exp = explain_box.text_frame
tf_exp.word_wrap = True

p_e_title = tf_exp.paragraphs[0]
p_e_title.text = "The Reality of 'Hardware Eye'"
p_e_title.font.name = FONT_HEADING
p_e_title.font.size = Pt(20)
p_e_title.font.bold = True
p_e_title.font.color.rgb = TEXT_WHITE
p_e_title.space_after = Pt(20)

bullets_exp = [
    ("On-Device Print", "Generated directly on the Arduino Nano using our visualDebugPrint() helper, which outputs a 24x24 ASCII grid downsampled from the 96x96 grayscale buffer."),
    ("Contrast & Scale", "Highlights the blocky, low-contrast grayscale crop that the model is forced to interpret."),
    ("Explainable Failure", "Visualizing this input makes it immediately clear why classification fails: color, edge, and texture details are completely washed away, leaving only generic shapes.")
]
for title, text in bullets_exp:
    p = tf_exp.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

# ===========================================================================
# SLIDE 8 (NEW): Real Snapshot: Reality vs. Camera View
# ===========================================================================
slide8_new = prs.slides.add_slide(blank_layout)
set_slide_background(slide8_new)
add_header(slide8_new, "Real Snapshot: Reality vs. Camera View")

# Image on the left (4:3 aspect ratio fitted inside card)
create_card(slide8_new, Inches(0.8), Inches(1.8), Inches(6.4), Inches(4.8))

# Fit the image. Original is 4080x3072. We fit it to height 4.5" -> width 6.0"
# Center it inside the 6.4" x 4.8" card
slide8_new.shapes.add_picture(
    'present/comparison of reality and what camera see.jpg',
    left=Inches(1.0),
    top=Inches(1.95),
    width=Inches(6.0),
    height=Inches(4.5)
)

# Text description card on the right
create_card(slide8_new, Inches(7.6), Inches(1.8), Inches(4.9), Inches(4.8))
explain_snapshot_box = slide8_new.shapes.add_textbox(Inches(7.9), Inches(2.1), Inches(4.3), Inches(4.2))
tf_snap = explain_snapshot_box.text_frame
tf_snap.word_wrap = True

p_snap_title = tf_snap.paragraphs[0]
p_snap_title.text = "Visual Gap in Practice"
p_snap_title.font.name = FONT_HEADING
p_snap_title.font.size = Pt(20)
p_snap_title.font.bold = True
p_snap_title.font.color.rgb = TEXT_WHITE
p_snap_title.space_after = Pt(20)

bullets_snap = [
    ("Human Perspective", "Full-resolution, high-contrast, color-rich image where food details and boundaries are distinct."),
    ("Microcontroller View", "Grayscale, 96x96 resolution downsampled image where colors merge and edges blur."),
    ("Why It Matters", "Differentiating complex dishes like pizza vs. cheesecake requires color and texture. When reduced to a blocky silhouette, the classification boundaries overlap completely, leading to high-uncertainty predictions.")
]
for title, text in bullets_snap:
    p = tf_snap.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = ACCENT_TEAL
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

# ===========================================================================
# SLIDE 9 (NEW): Transition Slide: DEMO TIME!
# ===========================================================================
slide9_new = prs.slides.add_slide(blank_layout)
set_slide_background(slide9_new)

# Giant centered text box
demo_box = slide9_new.shapes.add_textbox(Inches(1.0), Inches(2.2), Inches(11.333), Inches(4.0))
tf_demo = demo_box.text_frame
tf_demo.word_wrap = True

p_demo = tf_demo.paragraphs[0]
p_demo.text = "DEMO TIME!"
p_demo.font.name = FONT_HEADING
p_demo.font.size = Pt(64)
p_demo.font.bold = True
p_demo.font.color.rgb = ACCENT_TEAL
p_demo.alignment = PP_ALIGN.CENTER
p_demo.space_after = Pt(20)

p_demo_sub = tf_demo.add_paragraph()
p_demo_sub.text = "Let's watch the real-time pantry logging in action..."
p_demo_sub.font.name = FONT_BODY
p_demo_sub.font.size = Pt(18)
p_demo_sub.font.color.rgb = TEXT_MUTED
p_demo_sub.alignment = PP_ALIGN.CENTER

# ===========================================================================
# SLIDE 10: Key Takeaways & Lessons Learned
# ===========================================================================
slide10 = prs.slides.add_slide(blank_layout)
set_slide_background(slide10)
add_header(slide10, "Lessons Learned & Edge CV Realities")

create_card(slide10, Inches(0.8), Inches(1.8), Inches(5.6), Inches(4.8))
lessons_box = slide10.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.0), Inches(4.2))
tf_less = lessons_box.text_frame
tf_less.word_wrap = True

p_l_title = tf_less.paragraphs[0]
p_l_title.text = "Key Takeaways"
p_l_title.font.name = FONT_HEADING
p_l_title.font.size = Pt(20)
p_l_title.font.bold = True
p_l_title.font.color.rgb = ACCENT_TEAL
p_l_title.space_after = Pt(15)

bullets_less = [
    ("RAM is the Hard Constraint", "In TinyML, memory footprint of weights (Flash) is secondary. Activation tensors (RAM) determine model feasibility."),
    ("Sufficient Feature Size", "Vision classification requires minimum information capacity. Grayscale 96x96 images lack the entropy required to differentiate complex dishes."),
    ("Constraint-Driven Scope", "CV on sub-256KB RAM devices is only practical for high-contrast, geometric shapes under rigid lighting (e.g., bar codes, simple gestures).")
]
for title, text in bullets_less:
    p = tf_less.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

create_card(slide10, Inches(6.9), Inches(1.8), Inches(5.6), Inches(4.8))
future_box = slide10.shapes.add_textbox(Inches(7.2), Inches(2.1), Inches(5.0), Inches(4.2))
tf_fut = future_box.text_frame
tf_fut.word_wrap = True

p_f_title = tf_fut.paragraphs[0]
p_f_title.text = "Future Directions"
p_f_title.font.name = FONT_HEADING
p_f_title.font.size = Pt(20)
p_f_title.font.bold = True
p_f_title.font.color.rgb = ACCENT_ORANGE
p_f_title.space_after = Pt(15)

bullets_fut = [
    ("Hardware Upgrades", "Select microcontrollers with external PSRAM (e.g., ESP32-S3) to support larger activation buffers without scaling down resolution."),
    ("Sensor Fusion", "Integrate low-power non-vision sensors, such as gas/volatile compound sensors, weight cells, or temperature sensors to track freshness directly."),
    ("Custom In-Situ Data", "Collect a custom in-situ dataset inside the fridge itself rather than relying on noisy web-scraped databases like Food-101.")
]
for title, text in bullets_fut:
    p = tf_fut.add_paragraph()
    p.text = f"\u2022  {title}: "
    p.font.bold = True
    p.font.size = Pt(14)
    p.font.color.rgb = TEXT_WHITE
    p.space_after = Pt(2)
    
    run = p.add_run()
    run.text = text
    run.font.bold = False
    run.font.color.rgb = TEXT_MUTED
    p.space_after = Pt(12)

# Save the presentation directly in the 'present' directory
try:
    prs.save("present/pantry_presentation.pptx")
    print("Presentation updated and saved successfully as present/pantry_presentation.pptx")
except PermissionError:
    prs.save("present/pantry_presentation_updated.pptx")
    print("WARNING: present/pantry_presentation.pptx was locked (likely open in PowerPoint). Saved copy to present/pantry_presentation_updated.pptx instead!")
